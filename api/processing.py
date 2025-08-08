from pathlib import Path
import fitz  # pymupdf
from docx import Document
from pptx import Presentation
from gtts import gTTS

def extract_text_from_pdf(path: Path) -> str:
    doc = fitz.open(path)
    text = []
    for page in doc:
        try:
            text.append(page.get_text())
        except Exception:
            pass
    return '\n'.join(text)

def extract_text_from_docx(path: Path) -> str:
    doc = Document(path)
    paras = [p.text for p in doc.paragraphs]
    return '\n'.join(paras)

def extract_text_from_pptx(path: Path) -> str:
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slides_text.append(shape.text)
    return '\n'.join(slides_text)

def make_audio(text: str, out_path: Path):
    # gTTS has limits; we clip to safe length
    t = gTTS(text[:4000])
    t.save(str(out_path))

def save_html(html: str, out_path: Path):
    html_doc = f"<!doctype html><html><head><meta charset=\"utf-8\"><meta name=\"viewport\" content=\"width=device-width,initial-scale=1\"><title>Accessible version</title></head><body>{html}</body></html>"
    out_path.write_text(html_doc, encoding='utf-8')